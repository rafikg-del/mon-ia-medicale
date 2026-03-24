import hashlib
import json
import os
import re
import tempfile
import base64

import streamlit as st
from dotenv import load_dotenv

load_dotenv()

# ── CONFIG ────────────────────────────────────────────────────────────────────
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
XAI_API_KEY    = os.getenv("XAI_API_KEY")

GROK_MODEL = "grok-3"

GEMINI_FALLBACK = [
    "gemini-2.0-flash",
    "gemini-1.5-flash",
    "gemini-1.5-flash-latest",
    "gemini-1.5-pro",
    "gemini-1.5-flash-8b",
    "gemini-2.0-flash-lite",
]

CHAT_CTX_CHARS = 60_000
MAX_DOC_CHARS  = 150_000
MAX_SLIDES_IMAGES = 70  # max slides à envoyer en images

# ── PROMPTS ───────────────────────────────────────────────────────────────────

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# AGENT 1 — GEMINI : EXTRACTION EXHAUSTIVE (TEXTE)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
EXTRACTION_PROMPT_TEXT = """\
Tu es un agent d'extraction médicale de niveau expert pour le DIUE MAPS \
(Micronutrition, Alimentation, Prévention et Santé).

MISSION : Extraire la TOTALITÉ du contenu de ce document de cours sans aucune perte d'information.

RÈGLES ABSOLUES :
1. EXHAUSTIVITÉ TOTALE : chaque fait, chiffre, mécanisme, exemple, digression du professeur, \
question/réponse d'étudiant, référence à un schéma ou slide doit être extrait.
2. ZÉRO INVENTION : n'ajoute RIEN qui ne soit pas dans le document.
3. PRÉSERVER LES NUANCES : si le professeur dit "on n'a pas prouvé" ou "c'est discuté", \
retranscris cette nuance exactement.
4. DISTINGUER les types de contenu :
   - [COURS] = contenu enseigné comme fait établi
   - [DIGRESSION_PROF] = opinion personnelle, hypothèse, anecdote du professeur
   - [QUESTION_ETUDIANT] = question posée par un étudiant et réponse du prof
   - [SCHEMA] = description d'un schéma, slide, ou visuel commenté

Retourne un JSON structuré avec :
{{
  "titre": "thème principal du cours",
  "professeur": "nom du professeur si mentionné",
  "contexte": "module, date, formation (DIUE MAPS, etc.)",
  "type_source": "transcript | diapo_texte | diapo_images | mixte",
  "alerte_source_pauvre": false,
  "concepts_fondamentaux": [
    {{
      "concept": "nom du concept",
      "type": "COURS | DIGRESSION_PROF",
      "contenu_complet": "explication INTÉGRALE sans résumé",
      "cascade_causale": "A → B → C si applicable",
      "exemples_cites": ["tous les exemples mentionnés par le prof"],
      "nuances_prof": "avertissements, limites, doutes exprimés"
    }}
  ],
  "schemas_et_slides": [
    {{
      "reference": "slide X ou description du schéma",
      "contenu": "ce que le schéma montre",
      "message_implicite": "interprétation du visuel"
    }}
  ],
  "experiences_citees": [
    {{
      "nom": "nom de l'expérience ou étude",
      "protocole": "description",
      "resultat": "résultat",
      "conclusion_prof": "ce qu'en dit le professeur"
    }}
  ],
  "digressions_professeur": [
    {{
      "sujet": "thème de la digression",
      "contenu_integral": "TOUT ce que le prof a dit, sans couper"
    }}
  ],
  "questions_etudiants": [
    {{
      "question": "question posée",
      "reponse": "réponse du professeur"
    }}
  ],
  "cofacteurs_metabolites_cites": ["liste de tous les cofacteurs, vitamines, minéraux, \
métabolites mentionnés dans le cours"],
  "liens_micronutrition": ["tout ce que le cours dit sur le lien avec la micronutrition, \
l'alimentation, la supplémentation"],
  "points_critiques": ["complications, urgences, contre-indications MENTIONNÉES dans le cours"],
  "avertissements_prof": ["mises en garde explicites du professeur"]
}}

DOCUMENT :
---
{document_text}
---

Retourne UNIQUEMENT le JSON. Sois EXHAUSTIF. Ne résume JAMAIS.
"""

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# AGENT 1 — GEMINI : EXTRACTION MULTIMODALE (SLIDES = IMAGES)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
EXTRACTION_PROMPT_VISUAL = """\
Tu es un professeur expert en médecine fonctionnelle et biochimie, spécialisé DIUE MAPS \
(Micronutrition, Alimentation, Prévention et Santé).

MISSION : Tu reçois les IMAGES des slides d'un diaporama de cours. \
Tu dois les analyser comme le ferait le professeur qui les commente à l'oral.

RÈGLES ABSOLUES :
1. DÉCRIS TOUT ce que tu vois sur chaque slide : texte, schémas, flèches, cascades, \
   formules chimiques, structures moléculaires, tableaux, graphiques, légendes, couleurs \
   significatives, hiérarchie visuelle.
2. INTERPRÈTE chaque schéma : que montre-t-il ? quelle cascade mécanistique ? \
   quel message implicite ? quelles flèches causales ?
3. RESTE EXCLUSIVEMENT sur ce que montrent les images. N'invente RIEN qui ne soit pas \
   visible sur les slides. Si un schéma est ambigu, signale l'ambiguïté.
4. RECONSTRUIT la logique pédagogique : dans quel ordre les slides s'enchaînent, \
   quelle progression conceptuelle, quel fil conducteur.
5. Pour chaque slide, identifie :
   - Le CONTENU FACTUEL (texte visible, données, formules)
   - Le MESSAGE DU SCHÉMA (ce que le visuel communique au-delà du texte)
   - Les CASCADES CAUSALES implicites (flèches, enchaînements)
   - Les COFACTEURS / MÉTABOLITES visibles (vitamines, ions, molécules)
   - Les LIENS MICRONUTRITION implicites ou explicites

Retourne un JSON structuré avec :
{{
  "titre": "thème du diaporama",
  "professeur": "nom si visible",
  "contexte": "module, formation",
  "type_source": "diapo_images",
  "alerte_source_pauvre": false,
  "nombre_slides": N,
  "slides": [
    {{
      "numero": 1,
      "titre_slide": "titre visible ou déduit",
      "texte_visible": "tout le texte lisible sur la slide",
      "description_schemas": "description détaillée de tout schéma/visuel",
      "cascades_causales": "A → B → C identifiées dans les flèches/schémas",
      "formules_chimiques": "toute formule ou structure visible",
      "cofacteurs_visibles": ["métabolites, vitamines, ions visibles"],
      "message_pedagogique": "ce que cette slide enseigne, quel concept elle illustre",
      "liens_avec_slides_precedentes": "comment elle s'inscrit dans la progression"
    }}
  ],
  "synthese_globale": {{
    "concepts_fondamentaux": [
      {{
        "concept": "nom",
        "contenu_complet": "explication reconstituée à partir des visuels",
        "cascade_causale": "A → B → C",
        "cofacteurs_impliques": ["liste"]
      }}
    ],
    "cofacteurs_metabolites_identifies": ["liste globale"],
    "liens_micronutrition": ["tous les liens identifiés"]
  }}
}}

Analyse TOUTES les slides ci-dessous. Sois EXHAUSTIF sur les visuels.
"""

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# AGENT 2 — GROK : SYNTHÈSE FICHE RECONSTRUITE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SYNTHESIS_PROMPT = """\
Tu es un expert en médecine fonctionnelle, biochimie et pédagogie avancée \
(niveau médecin senior / chercheur), spécialisé dans le DIUE MAPS \
(Micronutrition, Alimentation, Prévention et Santé).

Ta mission : transformer les données extraites ci-dessous en une FICHE DE COURS RECONSTRUITE, \
exhaustive et directement actionnable en pratique de micronutrition.

━━━ DONNÉES EXTRAITES ━━━
{extracted_data}

━━━ RÈGLES ABSOLUES ━━━

1. ❌ ZÉRO CONTENU INVENTÉ : n'ajoute AUCUN traitement, posologie, examen complémentaire, \
pronostic ou espérance de vie qui ne figure pas dans les données extraites. \
Si tu reconstruis un lien causal implicite, signale-le avec [Reconstruction].

2. ❌ ZÉRO PERTE D'INFORMATION : chaque concept, exemple, digression, expérience, \
nuance et avertissement des données extraites DOIT apparaître dans ta fiche.

3. ❌ NE PAS utiliser un template de pathologie clinique (pas de "diagnostic différentiel", \
"traitement symptomatique", "pronostic" SAUF si le professeur en parle explicitement).

4. ✅ RECONSTRUIRE la logique causale : pour chaque concept, expliciter la cascade \
mécanistique (cause → mécanisme → conséquence → intervention possible).

5. ✅ DISTINGUER les niveaux de preuve :
   - **[Cours]** = enseigné comme fait établi
   - **(Digression Prof)** = opinion/hypothèse personnelle du professeur
   - **[Reconstruction]** = lien causal que TU reconstruis à partir du cours
   - **[Visuel]** = information extraite d'un schéma/slide (pas du texte oral)

6. ✅ ORIENTER vers l'action micronutritionnelle : chaque section doit mener vers \
un bottleneck identifiable, un biomarqueur mesurable, ou une intervention possible \
en médecine fonctionnelle / micronutrition.

7. ✅ Si la source est un DIAPORAMA (slides), reconstruire le commentaire du professeur \
tel qu'il l'aurait fait à l'oral : expliquer les schémas, interpréter les flèches, \
expliciter les cascades causales implicites. Signaler avec **[Visuel]**.

━━━ FORMAT DE SORTIE OBLIGATOIRE ━━━

# 🧠 {title}

> **Résumé** : 2-3 phrases de cadrage (contexte DIUE MAPS).

---

## F1 — [Premier grand thème du cours]

### Définitions & cadre
• Définitions telles qu'enseignées ou visibles sur les slides
• Contexte clinique / enjeux en micronutrition

### Physiopathologie CAUSALE
• Cascade mécanistique complète : A → B → C → D
• Boucles de régulation
• Interactions inter-systèmes
*(Ne JAMAIS rester descriptif — toujours reconstruire la logique)*

### Apports des slides/schémas
• **[Visuel]** Description et interprétation de chaque schéma pertinent
• Messages implicites des visuels (flèches, couleurs, hiérarchie)

### Exemples & expériences cités
• [Reproduire INTÉGRALEMENT chaque exemple visible ou mentionné]

### Digressions du professeur
• *(Digression Prof)* [Contenu intégral si présent dans la source]

---

*(Répéter pour F2, F3... autant de fiches que de grands thèmes)*

---

## 🔬 BOTTLENECKS BIOLOGIQUES

| Bottleneck | Mécanisme | Conséquence | Intervention micronutritionnelle |
|---|---|---|---|

## 🧠 ALGORITHMES DÉCISIONNELS

➤ **IF** [condition biologique mesurable] **THEN** [intervention adaptée]

## 🗺️ MAPPING CORTEX

**Biomarqueur** → **Voie métabolique** → **Mécanisme** → **Intervention**

## ⚠️ POINTS DE VIGILANCE

## 💡 POINTS CLÉS / MNÉMOTECHNIQUES

---

━━━ VÉRIFICATION FINALE ━━━
☐ Chaque concept des données extraites apparaît dans la fiche
☐ Aucun traitement/posologie/pronostic n'a été inventé
☐ Chaque digression du prof est conservée et signalée
☐ Chaque schéma est décrit et interprété [Visuel]
☐ Les algorithmes IF/THEN sont orientés micronutrition
☐ Les bottlenecks identifient des cofacteurs/métabolites actionnables

Rédige la fiche COMPLÈTE maintenant.
"""

GROK_SYSTEM = """\
Tu es un assistant expert en médecine fonctionnelle et micronutrition pour le DIUE MAPS.

{fiche_block}

{context_block}

RÈGLES :
- Appuie-toi PRIORITAIREMENT sur la fiche et le cours source pour répondre.
- Ne JAMAIS inventer de contenu non présent dans les sources.
- Si une question dépasse le contenu du cours, dis-le clairement avant de compléter \
avec tes connaissances (en le signalant).
- Oriente tes réponses vers la pratique de micronutrition / médecine fonctionnelle.
- Distingue toujours ce qui vient du cours [Cours] vs tes ajouts [Connaissance externe].
- Réponds en français.
"""

# ── DOCUMENT EXTRACTION ───────────────────────────────────────────────────────
def _file_hash(f) -> str:
    f.seek(0); h = hashlib.md5(f.read()).hexdigest(); f.seek(0)
    return h


def extract_document_text(uploaded_file, file_type: str) -> str:
    """Extraction TEXTE classique (transcripts, docx, etc.)"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
        tmp.write(uploaded_file.read())
        path = tmp.name
    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(path)
            out = "\n".join(f"--- Page {i+1} ---\n{p.get_text()}" for i, p in enumerate(doc))
            doc.close()
            return out
        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            out = ""
            for i, slide in enumerate(prs.slides):
                out += f"\n--- Slide {i+1} ---\n"
                if slide.shapes.title:
                    out += f"TITRE: {slide.shapes.title.text}\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                        out += shape.text + "\n"
            return out
        elif file_type == "docx":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return ""
    finally:
        os.unlink(path)


def extract_slides_as_images(uploaded_file, file_type: str) -> list[dict]:
    """Convertit PDF/PPTX en images pour extraction visuelle multimodale.
    Retourne une liste de {"page": int, "mime": str, "data_b64": str}."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
        tmp.write(uploaded_file.read())
        path = tmp.name

    images = []
    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(path)
            for i, page in enumerate(doc):
                if i >= MAX_SLIDES_IMAGES:
                    break
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("jpeg")
                images.append({
                    "page": i + 1,
                    "mime": "image/jpeg",
                    "data_b64": base64.b64encode(img_bytes).decode()
                })
            doc.close()

        elif file_type == "pptx":
            # Convertir PPTX → PDF via LibreOffice, puis PDF → images
            import subprocess
            out_dir = tempfile.mkdtemp()
            try:
                subprocess.run([
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", out_dir, path
                ], check=True, capture_output=True, timeout=120)
                pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(path))[0] + ".pdf")
                if os.path.exists(pdf_path):
                    import fitz
                    doc = fitz.open(pdf_path)
                    for i, page in enumerate(doc):
                        if i >= MAX_SLIDES_IMAGES:
                            break
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("jpeg")
                        images.append({
                            "page": i + 1,
                            "mime": "image/jpeg",
                            "data_b64": base64.b64encode(img_bytes).decode()
                        })
                    doc.close()
                    os.unlink(pdf_path)
            except (subprocess.CalledProcessError, FileNotFoundError):
                # LibreOffice non disponible — fallback texte seul
                pass
            finally:
                try:
                    os.rmdir(out_dir)
                except OSError:
                    pass
    finally:
        os.unlink(path)

    return images


def _is_slide_document(file_name: str, text: str) -> bool:
    """Heuristique : détecte si c'est un diaporama (peu de texte, beaucoup de slides)."""
    name_lower = file_name.lower()
    if "diapo" in name_lower or "slide" in name_lower or "ppt" in name_lower:
        return True
    # PDF avec peu de texte par page = probablement des slides
    pages = text.count("--- Page ")
    if pages > 0:
        chars_per_page = len(text) / pages
        if chars_per_page < 500 and pages > 10:
            return True
    return False


# ── AGENTS ────────────────────────────────────────────────────────────────────
class GeminiAgent:
    def __init__(self):
        import google.generativeai as genai
        self._genai = genai
        genai.configure(api_key=GOOGLE_API_KEY)

        try:
            available = [m.name for m in genai.list_models()
                         if "generateContent" in m.supported_generation_methods]
        except Exception:
            available = []

        if not available:
            raise RuntimeError("Aucun modèle Gemini accessible avec cette clé API.")

        self.model = None
        self.model_name = None
        for name in GEMINI_FALLBACK:
            normalized = name if name.startswith("models/") else f"models/{name}"
            if normalized not in available:
                continue
            try:
                m = genai.GenerativeModel(name)
                m.generate_content("test", generation_config={"max_output_tokens": 1, "temperature": 0})
                self.model = m
                self.model_name = name
                break
            except Exception:
                continue

        if self.model is None and available:
            fallback_name = available[0].removeprefix("models/")
            self.model = genai.GenerativeModel(fallback_name)
            self.model_name = fallback_name

        if self.model is None:
            raise RuntimeError(f"Aucun modèle Gemini disponible.")

    def extract_text(self, text: str) -> str:
        """Extraction depuis du texte pur (transcripts, docx)."""
        return self.model.generate_content(
            EXTRACTION_PROMPT_TEXT.format(document_text=text),
            generation_config={"max_output_tokens": 8192, "temperature": 0.1},
        ).text

    def extract_visual(self, slide_images: list[dict]) -> str:
        """Extraction multimodale : envoie les images des slides à Gemini."""
        parts = [EXTRACTION_PROMPT_VISUAL]
        for img in slide_images:
            parts.append({
                "mime_type": img["mime"],
                "data": base64.b64decode(img["data_b64"])
            })
            parts.append(f"\n[Slide {img['page']}]\n")

        return self.model.generate_content(
            parts,
            generation_config={"max_output_tokens": 8192, "temperature": 0.1},
        ).text

    def extract_combined(self, text: str, slide_images: list[dict]) -> str:
        """Extraction combinée : texte + images des slides."""
        # D'abord extraction visuelle des slides
        visual_data = self.extract_visual(slide_images) if slide_images else ""
        # Puis extraction texte
        text_data = self.extract_text(text) if text.strip() else ""

        if visual_data and text_data:
            # Fusionner les deux extractions
            merge_prompt = f"""\
Tu reçois deux extractions du MÊME cours :
1) Extraction VISUELLE (depuis les images des slides) :
{visual_data}

2) Extraction TEXTUELLE (depuis le texte/transcript) :
{text_data}

MISSION : Fusionne ces deux extractions en un SEUL JSON unifié.
- Conserve TOUTE l'information des deux sources sans perte.
- L'extraction visuelle apporte les schémas, flèches, structures moléculaires.
- L'extraction textuelle apporte les explications détaillées, digressions, exemples.
- En cas de redondance, garde la version la plus complète.
- En cas d'information complémentaire, fusionne.
Retourne le JSON fusionné.
"""
            return self.model.generate_content(
                merge_prompt,
                generation_config={"max_output_tokens": 8192, "temperature": 0.1},
            ).text
        return visual_data or text_data


class GrokAgent:
    def __init__(self):
        from openai import OpenAI
        self.client = OpenAI(api_key=XAI_API_KEY, base_url="https://api.x.ai/v1")
        self._messages: list[dict] = []
        self._system = ""

    def synthesize(self, data: str, title: str) -> str:
        res = self.client.chat.completions.create(
            model=GROK_MODEL,
            max_tokens=8192,
            messages=[{"role": "user", "content": SYNTHESIS_PROMPT.format(
                extracted_data=data, title=title,
            )}],
        )
        return res.choices[0].message.content

    def initialize_chat(self, context: str = "", fiche: str = "") -> None:
        self._messages = []
        self._system = GROK_SYSTEM.format(
            fiche_block=(
                f"=== FICHE MÉDICALE (référence prioritaire) ===\n{fiche}\n=== FIN ==="
                if fiche else "Aucune fiche générée."
            ),
            context_block=(
                f"=== COURS SOURCE ===\n{context[:CHAT_CTX_CHARS]}\n=== FIN ==="
                if context else ""
            ),
        )

    def chat(self, message: str, context: str = "", fiche: str = "") -> str:
        if not self._system:
            self.initialize_chat(context, fiche)
        self._messages.append({"role": "user", "content": message})
        res = self.client.chat.completions.create(
            model=GROK_MODEL,
            max_tokens=4096,
            messages=[{"role": "system", "content": self._system}] + self._messages,
        )
        reply = res.choices[0].message.content
        self._messages.append({"role": "assistant", "content": reply})
        return reply

# ── LATEX RENDERER ────────────────────────────────────────────────────────────
_DISPLAY_MATH = re.compile(r'\$\$(.+?)\$\$', re.DOTALL)

def render_latex(content: str) -> None:
    for i, part in enumerate(_DISPLAY_MATH.split(content)):
        if i % 2 == 0:
            if part.strip():
                st.markdown(part)
        else:
            st.latex(part.strip())

# ── SIDEBAR CHAT ──────────────────────────────────────────────────────────────
def _ctx_hash(t: str) -> str:
    return str(hash(t[:500])) if t else ""

def render_sidebar(doc_ctx: str = "", fiche: str = "") -> None:
    for k, v in [("grok", None), ("chat_msgs", []), ("chat_dh", ""), ("chat_fh", "")]:
        if k not in st.session_state:
            st.session_state[k] = v

    dh, fh = _ctx_hash(doc_ctx), _ctx_hash(fiche)
    if (dh != st.session_state.chat_dh or fh != st.session_state.chat_fh) and (doc_ctx or fiche):
        if st.session_state.grok is None:
            st.session_state.grok = GrokAgent()
        st.session_state.grok.initialize_chat(doc_ctx, fiche)
        st.session_state.chat_dh = dh
        st.session_state.chat_fh = fh

    status = "✅ Fiche + cours" if fiche else ("📄 Cours chargé" if doc_ctx else "Aucun contexte")

    with st.sidebar:
        st.markdown(f"""
        <div style="background:linear-gradient(135deg,#0a0a23,#1a0533);
                    padding:16px;border-radius:12px;margin-bottom:12px;">
            <div style="color:#fff;font-size:17px;font-weight:700;">✨ Grok-3 — Chat DIUE MAPS</div>
            <div style="color:rgba(255,255,255,0.5);font-size:11px;margin-top:4px;">{status}</div>
        </div>""", unsafe_allow_html=True)

        if not fiche and not doc_ctx:
            st.warning("Importez un document pour activer le chat.")
        elif not fiche:
            st.info("💡 Générez une fiche pour des réponses précises.")
        else:
            st.success("🎯 Grok-3 a accès à votre fiche.")

        st.divider()

        for msg in st.session_state.chat_msgs:
            if msg["role"] == "user":
                st.markdown(f"""<div style="background:#e3f2fd;border-radius:10px 10px 2px 10px;
                    padding:9px 12px;margin:5px 0;font-size:13px;color:#0d47a1;">
                    👤 <b>Vous</b><br>{msg['content']}</div>""", unsafe_allow_html=True)
            else:
                st.markdown("""<div style="background:#ede7f6;border-radius:10px 10px 10px 2px;
                    padding:9px 12px;margin:3px 0;font-size:13px;color:#311b92;">
                    ✨ <b>Grok-3</b></div>""", unsafe_allow_html=True)
                st.markdown(msg["content"])

        user_input = st.chat_input("Posez une question...")
        if user_input:
            if st.session_state.grok is None:
                st.session_state.grok = GrokAgent()
            st.session_state.chat_msgs.append({"role": "user", "content": user_input})
            with st.spinner("Grok-3 réfléchit..."):
                reply = st.session_state.grok.chat(user_input, doc_ctx, fiche)
            st.session_state.chat_msgs.append({"role": "assistant", "content": reply})
            st.rerun()

        if st.session_state.chat_msgs:
            st.divider()
            if st.button("🗑️ Effacer", use_container_width=True):
                st.session_state.chat_msgs = []
                st.session_state.grok = None
                st.session_state.chat_dh = ""
                st.session_state.chat_fh = ""
                st.rerun()

# ── APP ───────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="MedFiche AI — DIUE MAPS", page_icon="🏥", layout="wide")

st.markdown("""
<style>
.hdr{background:linear-gradient(135deg,#0a0a23,#1a0533,#001a33);
     padding:24px 30px;border-radius:14px;margin-bottom:22px;
     box-shadow:0 8px 30px rgba(0,0,0,0.4);}
.hdr h1{color:#fff;margin:0;font-size:26px;}
.hdr p{color:rgba(255,255,255,0.6);margin:8px 0 0;font-size:13px;}
.hdr .tags{display:flex;gap:8px;margin-top:12px;flex-wrap:wrap;}
.tag{padding:4px 11px;border-radius:16px;font-size:11px;font-weight:600;color:#fff;}
.t1{background:linear-gradient(90deg,#1a73e8,#0d47a1);}
.t2{background:linear-gradient(90deg,#7c4dff,#311b92);}
.card{border-radius:10px;padding:13px 17px;margin:7px 0;
      display:flex;align-items:center;gap:11px;font-size:14px;}
.c1{background:#e3f2fd;border-left:4px solid #1565c0;}
.c2{background:#ede7f6;border-left:4px solid #6a1b9a;}
.c3{background:#e8f5e9;border-left:4px solid #2e7d32;}
.badge{display:inline-flex;align-items:center;gap:7px;
       background:linear-gradient(90deg,#1a73e8,#7c4dff);
       color:#fff;padding:5px 14px;border-radius:18px;
       font-size:12px;font-weight:600;margin-bottom:16px;}
#MainMenu,footer{visibility:hidden;}
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="hdr">
  <h1>🏥 MedFiche AI — DIUE MAPS</h1>
  <p>Fiches de cours reconstruites · Extraction visuelle multimodale · Micronutrition</p>
  <div class="tags">
    <span class="tag t1">🔍 Gemini · Extraction texte + images</span>
    <span class="tag t2">✨ Grok-3 · Reconstruction causale</span>
  </div>
</div>""", unsafe_allow_html=True)

if not all([GOOGLE_API_KEY, XAI_API_KEY]):
    missing = [k for k, v in {"GOOGLE_API_KEY": GOOGLE_API_KEY, "XAI_API_KEY": XAI_API_KEY}.items() if not v]
    st.error(f"Clés API manquantes : {', '.join(missing)}")
    st.stop()

for k, v in [("doc_text", None), ("file_name", None), ("results", None), ("slide_images", None)]:
    if k not in st.session_state:
        st.session_state[k] = v

fiche_ctx = (st.session_state.results or {}).get("final_fiche", "") or ""
render_sidebar(st.session_state.doc_text or "", fiche_ctx)

tab1, tab2, tab3 = st.tabs(["📁 Import & Génération", "📋 Fiche Reconstruite", "🔬 Détails"])

# ── TAB 1 ─────────────────────────────────────────────────────────────────────
with tab1:
    uploaded = st.file_uploader("Cours (PDF, PPTX, DOCX)", type=["pdf", "pptx", "docx"])

    if uploaded:
        ftype = uploaded.name.split(".")[-1].lower()
        fh = _file_hash(uploaded)
        cache_key = f"doc_{fh}"
        if cache_key not in st.session_state:
            with st.spinner("Lecture du fichier..."):
                uploaded.seek(0)
                raw = extract_document_text(uploaded, ftype)

                # Extraire aussi les images si c'est un PDF ou PPTX
                slide_imgs = []
                if ftype in ("pdf", "pptx"):
                    uploaded.seek(0)
                    with st.spinner("Conversion slides → images..."):
                        slide_imgs = extract_slides_as_images(uploaded, ftype)

            for k in list(st.session_state.keys()):
                if k.startswith("doc_") and k != cache_key:
                    del st.session_state[k]
            st.session_state[cache_key]     = raw[:MAX_DOC_CHARS]
            st.session_state.doc_text       = st.session_state[cache_key]
            st.session_state.file_name      = uploaded.name
            st.session_state.slide_images   = slide_imgs
            st.session_state.results        = None

        text = st.session_state[cache_key]
        imgs = st.session_state.slide_images or []

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Caractères", f"{len(text):,}")
        c2.metric("Mots",       f"{len(text.split()):,}")
        c3.metric("Tokens (~)", f"{len(text)//4:,}")
        c4.metric("Slides (img)", f"{len(imgs)}")

        is_slides = _is_slide_document(st.session_state.file_name, text)
        if is_slides and imgs:
            st.info(f"📊 Document détecté comme **diaporama** ({len(imgs)} slides). "
                    f"Extraction visuelle multimodale activée.")
        elif is_slides and not imgs:
            st.warning("📊 Document détecté comme diaporama mais conversion images échouée. "
                       "Extraction texte seul (qualité réduite).")

        with st.expander("👁️ Aperçu texte", expanded=False):
            st.text(text[:2000] + ("..." if len(text) > 2000 else ""))

        if imgs:
            with st.expander(f"🖼️ Aperçu slides ({len(imgs)} images)", expanded=False):
                cols = st.columns(min(4, len(imgs)))
                for i, img in enumerate(imgs[:8]):
                    with cols[i % 4]:
                        st.image(base64.b64decode(img["data_b64"]),
                                 caption=f"Slide {img['page']}", use_container_width=True)
                if len(imgs) > 8:
                    st.caption(f"... et {len(imgs) - 8} slides supplémentaires")

    if st.session_state.doc_text:
        col, _ = st.columns([2, 3])
        with col:
            go = st.button("🚀 Générer la Fiche Reconstruite", type="primary", use_container_width=True)

        if go:
            bar    = st.progress(0.0)
            status = st.empty()
            cards  = st.empty()

            text = st.session_state.doc_text
            imgs = st.session_state.slide_images or []
            is_slides = _is_slide_document(st.session_state.file_name, text)

            def _show_cards(active: int, use_visual: bool) -> None:
                if use_visual:
                    agents = [
                        ("c1", "🖼️", "Agent 1a — Gemini Vision", f"Analyse visuelle {len(imgs)} slides"),
                        ("c3", "🔍", "Agent 1b — Gemini Texte", "Extraction texte + fusion"),
                        ("c2", "✨", "Agent 2 — Grok-3", "Reconstruction causale + micronutrition"),
                    ]
                else:
                    agents = [
                        ("c1", "🔍", "Agent 1 — Gemini", "Extraction exhaustive"),
                        ("c2", "✨", "Agent 2 — Grok-3", "Reconstruction causale + micronutrition"),
                    ]
                html = ""
                for i, (cls, ico, lbl, sub) in enumerate(agents, 1):
                    s = "⏳ En cours..." if i == active else ("✅ Terminé" if i < active else "⌛ En attente")
                    op = "1" if i <= active else "0.4"
                    html += f'<div class="card {cls}" style="opacity:{op};"><span style="font-size:20px">{ico}</span><div><b>{lbl}</b><br><span style="font-size:12px;color:#666">{sub} — {s}</span></div></div>'
                cards.markdown(html, unsafe_allow_html=True)

            try:
                use_visual = is_slides and len(imgs) > 0
                gemini = GeminiAgent()

                if use_visual:
                    # Étape 1a : extraction visuelle
                    bar.progress(0.05)
                    status.info(f"🖼️ Agent 1a — Gemini analyse {len(imgs)} slides visuellement...")
                    _show_cards(1, True)
                    ext = gemini.extract_combined(text, imgs)
                    bar.progress(0.5)
                    _show_cards(3, True)
                else:
                    # Extraction texte classique
                    bar.progress(0.05)
                    status.info("🔍 Agent 1 — Gemini extrait le corpus...")
                    _show_cards(1, False)
                    ext = gemini.extract_text(text)
                    bar.progress(0.5)
                    _show_cards(2, False)

                # Titre
                title = st.session_state.file_name or "Cours médical"
                try:
                    m = re.search(r'\{.*\}', ext, re.DOTALL)
                    if m:
                        title = json.loads(m.group()).get("titre", title)
                except Exception:
                    pass

                # Agent 2 — Grok-3 synthèse
                status.info("✨ Agent 2 — Grok-3 reconstruit la fiche causale...")
                _show_cards(3 if use_visual else 2, use_visual)
                final = GrokAgent().synthesize(ext, title)

                bar.progress(1.0)
                _show_cards(4 if use_visual else 3, use_visual)
                st.session_state.results = {
                    "extracted_data": ext,
                    "final_fiche": final,
                    "title": title,
                    "extraction_mode": "visuelle+texte" if use_visual else "texte",
                }
                status.success(f"🎉 Terminé ! Mode : **{'Visuel + Texte' if use_visual else 'Texte'}**. "
                               f"Consultez l'onglet **Fiche Reconstruite**.")

            except Exception as e:
                status.error(f"❌ Erreur : {e}")

# ── TAB 2 ─────────────────────────────────────────────────────────────────────
with tab2:
    if not st.session_state.results:
        st.info("👆 Importez un document et lancez la génération.")
    else:
        res = st.session_state.results
        col_t, col_dl = st.columns([5, 1])
        with col_t:
            mode = res.get("extraction_mode", "texte")
            st.markdown(f"## 📋 {res['title']}")
        with col_dl:
            st.download_button("⬇️ .md", data=res["final_fiche"],
                file_name=f"fiche_{res['title']}.md", mime="text/markdown",
                use_container_width=True)
        mode_label = "🖼️ Visuel+Texte" if mode == "visuelle+texte" else "📄 Texte"
        st.markdown(f'<div class="badge">{mode_label} · Gemini → Grok-3</div>', unsafe_allow_html=True)
        st.divider()
        render_latex(res["final_fiche"])

# ── TAB 3 ─────────────────────────────────────────────────────────────────────
with tab3:
    if not st.session_state.results:
        st.info("Lancez d'abord la génération.")
    else:
        res = st.session_state.results
        with st.expander("🔍 Extraction Gemini (JSON brut)", expanded=False):
            st.markdown(res["extracted_data"])
        with st.expander("✨ Fiche Reconstruite Grok-3", expanded=True):
            render_latex(res["final_fiche"])
